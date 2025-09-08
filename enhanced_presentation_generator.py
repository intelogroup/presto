#!/usr/bin/env python3
"""
Enhanced Presentation Generator with Template Design Matching

This script creates presentations that closely replicate the original template's
design using extracted design properties including colors, fonts, layouts, and styling.

Author: AI Assistant
Date: 2024
"""

import json
import os
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from jinja2 import Environment, FileSystemLoader, Template
from lxml import etree

from template_config import TemplateConfig, ContentType, get_default_config


class EnhancedPresentationGenerator:
    """
    Enhanced presentation generator that uses extracted design properties
    to create presentations matching the original template design.
    """
    
    def __init__(self, template_path: str, design_config_path: str = "template_analysis.json"):
        """
        Initialize the enhanced presentation generator.
        
        Args:
            template_path: Path to the PowerPoint template file (.pptx)
            design_config_path: Path to the enhanced design configuration JSON
        """
        self.template_path = Path(template_path)
        self.design_config_path = Path(design_config_path)
        self.jinja_env = Environment(loader=FileSystemLoader("templates/"))
        
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")
        
        # Load design configuration
        self.design_config = self._load_design_config()
        
        print(f"✅ Loaded design configuration with {len(self.design_config.get('color_scheme', {}))} colors")
        print(f"✅ Font configuration: {self.design_config.get('fonts', {}).get('primary', 'Default')}")
    
    def _load_design_config(self) -> Dict[str, Any]:
        """
        Load the enhanced design configuration.
        
        Returns:
            Dictionary containing design properties
        """
        try:
            with open(self.design_config_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except FileNotFoundError:
            print(f"⚠️  Design config not found: {self.design_config_path}")
            return self._get_default_design_config()
    
    def _get_default_design_config(self) -> Dict[str, Any]:
        """
        Get default design configuration if file is not found.
        
        Returns:
            Default design configuration
        """
        return {
            "slide_dimensions": {"width": 10, "height": 7.5},
            "color_scheme": {
                "primary": "#4472C4",
                "secondary": "#E7E6E6",
                "accent": "#70AD47",
                "text_dark": "#000000",
                "text_light": "#FFFFFF"
            },
            "fonts": {
                "primary": "Calibri",
                "title_size": 44,
                "heading_size": 32,
                "body_size": 18
            }
        }
    
    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """
        Convert hex color to RGB tuple.
        
        Args:
            hex_color: Hex color string (e.g., '#FF0000')
            
        Returns:
            RGB tuple (r, g, b)
        """
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def _apply_color_scheme(self, shape, color_type: str = "primary"):
        """
        Apply color scheme to a shape.
        
        Args:
            shape: python-pptx shape object
            color_type: Type of color to apply (primary, secondary, accent, etc.)
        """
        try:
            color_scheme = self.design_config.get("color_scheme", {})
            hex_color = color_scheme.get(color_type, "#4472C4")
            r, g, b = self._hex_to_rgb(hex_color)
            
            if hasattr(shape, 'fill'):
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(r, g, b)
        except Exception as e:
            print(f"⚠️  Could not apply color scheme: {e}")
    
    def _apply_font_styling(self, text_frame, font_type: str = "body"):
        """
        Apply font styling based on design configuration.
        
        Args:
            text_frame: python-pptx text frame object
            font_type: Type of font styling (title, heading, body)
        """
        try:
            fonts_config = self.design_config.get("fonts", {})
            font_family = fonts_config.get("primary", "Calibri")
            
            # Determine font size based on type
            size_map = {
                "title": fonts_config.get("title_size", 44),
                "heading": fonts_config.get("heading_size", 32),
                "subheading": fonts_config.get("subheading_size", 24),
                "body": fonts_config.get("body_size", 18),
                "caption": fonts_config.get("caption_size", 14)
            }
            
            font_size = size_map.get(font_type, 18)
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_family
                    run.font.size = Pt(font_size)
                    
                    # Apply text color
                    color_scheme = self.design_config.get("color_scheme", {})
                    text_color = color_scheme.get("text_dark", "#000000")
                    r, g, b = self._hex_to_rgb(text_color)
                    run.font.color.rgb = RGBColor(r, g, b)
                    
        except Exception as e:
            print(f"⚠️  Could not apply font styling: {e}")
    
    def _create_styled_slide(self, prs: Presentation, layout_index: int = 0) -> Any:
        """
        Create a new slide with enhanced styling.
        
        Args:
            prs: Presentation object
            layout_index: Index of the slide layout to use
            
        Returns:
            Styled slide object
        """
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply background styling if needed
        try:
            color_scheme = self.design_config.get("color_scheme", {})
            bg_color = color_scheme.get("background", None)
            
            if bg_color and bg_color != "#FFFFFF":
                # Apply background color
                background = slide.background
                fill = background.fill
                fill.solid()
                r, g, b = self._hex_to_rgb(bg_color)
                fill.fore_color.rgb = RGBColor(r, g, b)
        except Exception as e:
            print(f"⚠️  Could not apply background styling: {e}")
        
        return slide
    
    def load_data(self, data_path: str) -> Dict[str, Any]:
        """
        Load JSON data for template processing.
        
        Args:
            data_path: Path to JSON data file
            
        Returns:
            Dictionary containing the loaded data
        """
        with open(data_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    
    def process_text_with_jinja(self, text: str, data: Dict[str, Any]) -> str:
        """
        Process text content using Jinja2 templating.
        
        Args:
            text: Text content with Jinja2 template syntax
            data: Data dictionary for template rendering
            
        Returns:
            Processed text with variables replaced
        """
        template = self.jinja_env.from_string(text)
        return template.render(**data)
    
    def create_enhanced_title_slide(self, prs: Presentation, data: Dict[str, Any]) -> Any:
        """
        Create an enhanced title slide matching template design.
        
        Args:
            prs: Presentation object
            data: Data dictionary
            
        Returns:
            Created slide object
        """
        slide = self._create_styled_slide(prs, 0)  # Title slide layout
        
        # Update title
        if slide.shapes.title:
            title_text = data.get('presentation_data', {}).get('title_slide', {}).get('title', 'Presentation Title')
            processed_title = self.process_text_with_jinja(title_text, data)
            slide.shapes.title.text = processed_title
            self._apply_font_styling(slide.shapes.title.text_frame, "title")
        
        # Update subtitle
        try:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_text = data.get('presentation_data', {}).get('title_slide', {}).get('subtitle', 'Subtitle')
            processed_subtitle = self.process_text_with_jinja(subtitle_text, data)
            subtitle_placeholder.text = processed_subtitle
            self._apply_font_styling(subtitle_placeholder.text_frame, "heading")
        except (KeyError, IndexError):
            print("⚠️  Could not find subtitle placeholder")
        
        return slide
    
    def create_enhanced_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], data: Dict[str, Any]) -> Any:
        """
        Create an enhanced content slide matching template design.
        
        Args:
            prs: Presentation object
            slide_data: Data for this specific slide
            data: Full data dictionary
            
        Returns:
            Created slide object
        """
        slide = self._create_styled_slide(prs, 1)  # Content slide layout
        
        # Update title
        if slide.shapes.title:
            title_text = slide_data.get('title', 'Slide Title')
            processed_title = self.process_text_with_jinja(title_text, data)
            slide.shapes.title.text = processed_title
            self._apply_font_styling(slide.shapes.title.text_frame, "heading")
        
        # Update content
        try:
            content_placeholder = slide.placeholders[1]
            content = slide_data.get('content', [])
            
            if hasattr(content_placeholder, 'text_frame'):
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                if isinstance(content, list):
                    # Handle bullet points
                    for i, item in enumerate(content):
                        processed_item = self.process_text_with_jinja(str(item), data)
                        if i == 0:
                            text_frame.text = processed_item
                        else:
                            p = text_frame.add_paragraph()
                            p.text = processed_item
                            p.level = 0
                else:
                    # Handle single text content
                    processed_content = self.process_text_with_jinja(str(content), data)
                    text_frame.text = processed_content
                
                self._apply_font_styling(text_frame, "body")
                
        except (KeyError, IndexError) as e:
            print(f"⚠️  Could not update content placeholder: {e}")
        
        return slide

    def create_slide_from_xml(self, prs: Presentation, slide_xml: str, layout_name: str):
        """
        Create a slide from an XML string.

        Args:
            prs: Presentation object
            slide_xml: XML string representing the slide content
            layout_name: The name of the slide layout to use
        """
        try:
            # Find the layout by name
            layout = None
            for l in prs.slide_layouts:
                if l.name == layout_name:
                    layout = l
                    break
            
            if layout is None:
                print(f"⚠️  Layout '{layout_name}' not found. Using default.")
                layout = prs.slide_layouts[0]

            slide = prs.slides.add_slide(layout)
            
            parser = etree.XMLParser(recover=True)
            root = etree.fromstring(slide_xml.encode('utf-8'), parser=parser)

            for element in root.findall('.//shape'):
                shape_type = element.get('type')
                if shape_type == 'title':
                    if slide.shapes.title:
                        slide.shapes.title.text = element.text
                        self._apply_font_styling(slide.shapes.title.text_frame, "title")
                elif shape_type == 'body':
                    try:
                        body_shape = slide.placeholders[1]
                        body_shape.text = element.text
                        self._apply_font_styling(body_shape.text_frame, "body")
                    except (KeyError, IndexError):
                        print("⚠️  Could not find body placeholder for XML content")

        except Exception as e:
            print(f"❌ Error creating slide from XML: {e}")

    def generate_enhanced_presentation(self, data: Dict[str, Any], output_path: str) -> str:
        """
        Generate an enhanced presentation using extracted design properties.
        
        Args:
            data: Dictionary containing presentation data
            output_path: Path where the generated presentation will be saved
            
        Returns:
            Path to the generated presentation file
        """
        # Create a new presentation from the template
        prs = Presentation(self.template_path)
        
        # Set slide dimensions based on design config
        slide_dims = self.design_config.get("slide_dimensions", {})
        if "width" in slide_dims and "height" in slide_dims:
            try:
                # Convert to EMU if needed
                width = slide_dims["width"]
                height = slide_dims["height"]
                
                if width < 100000:  # Assume inches if small number
                    prs.slide_width = Inches(width)
                    prs.slide_height = Inches(height)
                else:  # Assume EMU if large number
                    prs.slide_width = width
                    prs.slide_height = height
                    
            except Exception as e:
                print(f"⚠️  Could not set slide dimensions: {e}")
        
        presentation_data = data.get('presentation_data', {})
        
        # Create title slide
        if 'title_slide' in presentation_data:
            template = self.jinja_env.get_template('title_slide.xml')
            slide_content = template.render(presentation_data['title_slide'])
            self.create_slide_from_xml(prs, slide_content, 'Title Slide')

        # Create content slides
        slide_count = 0
        for slide_info in presentation_data.get('slides', []):
            slide_type = slide_info.get('type', 'content')
            layout_name = ''
            template_name = ''

            if slide_type == 'title_and_content':
                layout_name = 'Title and Content'
                template_name = 'title_and_content_slide.xml'
            elif slide_type == 'section_header':
                layout_name = 'Section Header'
                template_name = 'section_header_slide.xml'
            elif slide_type == 'two_content':
                layout_name = 'Two Content'
                template_name = 'two_content_slide.xml'
            elif slide_type == 'comparison':
                layout_name = 'Comparison'
                template_name = 'comparison_slide.xml'
            elif slide_type == 'title_only':
                layout_name = 'Title Only'
                template_name = 'title_only_slide.xml'
            elif slide_type == 'blank':
                layout_name = 'Blank'
                template_name = 'blank_slide.xml'
            elif slide_type == 'content_with_caption':
                layout_name = 'Content with Caption'
                template_name = 'content_with_caption_slide.xml'
            elif slide_type == 'picture_with_caption':
                layout_name = 'Picture with Caption'
                template_name = 'picture_with_caption_slide.xml'


            if layout_name and template_name:
                template = self.jinja_env.get_template(template_name)
                slide_content = template.render(slide_info['data'])
                self.create_slide_from_xml(prs, slide_content, layout_name)
                slide_count += 1
        
        print(f"✅ Created {slide_count} enhanced content slides")
        
        # Save the generated presentation
        output_file = Path(output_path)
        prs.save(str(output_file))
        
        return str(output_file)


def main():
    """
    Main function to demonstrate the enhanced presentation generator.
    """
    # Configuration
    template_file = "Scientific Conference Slides.pptx"
    data_file = "presentation_content.json"
    output_file = "enhanced_generated_presentation.pptx"
    design_config_file = "enhanced_template_config.json"
    
    try:
        # Initialize the enhanced generator
        generator = EnhancedPresentationGenerator(template_file, design_config_file)
        
        # Load data
        data = generator.load_data(data_file)
        
        print(f"🚀 Starting enhanced presentation generation...")
        print(f"📄 Template: {template_file}")
        print(f"🎨 Design config: {design_config_file}")
        print(f"📊 Data source: {data_file}")
        print(f"🎯 Output: {output_file}")
        
        # Generate presentation
        result_path = generator.generate_enhanced_presentation(data, output_file)
        
        print(f"\n✅ Enhanced presentation generated successfully!")
        print(f"📁 Output file: {result_path}")
        print(f"🎨 Applied design properties from template analysis")
        print(f"📊 Processed {len(data.get('presentation_data', {}))} data sections")
        
    except FileNotFoundError as e:
        print(f"❌ File Error: {e}")
        print("💡 Make sure the template file and data file exist in the current directory")
    except json.JSONDecodeError as e:
        print(f"❌ JSON Error: {e}")
        print("💡 Check the JSON syntax in your data file")
    except Exception as e:
        print(f"❌ Unexpected error: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()