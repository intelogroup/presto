#!/usr/bin/env python3
"""
PowerPoint to PNG Converter
Converts PPTX files to PNG images using python-pptx and Pillow
"""

import sys
import os
from pathlib import Path
try:
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    import io
except ImportError as e:
    print(f"Missing required packages. Please install with:")
    print("pip install python-pptx pillow")
    sys.exit(1)

def convert_pptx_to_images(pptx_path, output_dir=None):
    """
    Convert PPTX file to PNG images
    
    Args:
        pptx_path (str): Path to the PPTX file
        output_dir (str): Output directory for PNG files
    """
    pptx_path = Path(pptx_path)
    
    if not pptx_path.exists():
        print(f"Error: File '{pptx_path}' not found")
        return False
    
    if output_dir is None:
        output_dir = pptx_path.parent / f"{pptx_path.stem}_images"
    else:
        output_dir = Path(output_dir)
    
    # Create output directory
    output_dir.mkdir(exist_ok=True)
    
    try:
        # Load presentation
        prs = Presentation(pptx_path)
        print(f"Loading presentation: {pptx_path.name}")
        print(f"Found {len(prs.slides)} slides")
        
        # Convert each slide
        for i, slide in enumerate(prs.slides, 1):
            print(f"Converting slide {i}...")
            
            # Create a basic image representation
            # Note: This is a simplified conversion as python-pptx doesn't directly support image export
            img = Image.new('RGB', (1920, 1080), 'white')
            draw = ImageDraw.Draw(img)
            
            # Try to get slide title and content
            slide_title = ""
            slide_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    if shape.text.strip() and not slide_title:
                        slide_title = shape.text.strip()
                    else:
                        slide_content.append(shape.text.strip())
            
            # Draw title
            if slide_title:
                try:
                    font = ImageFont.truetype("arial.ttf", 48)
                except:
                    font = ImageFont.load_default()
                
                # Calculate text position
                bbox = draw.textbbox((0, 0), slide_title, font=font)
                text_width = bbox[2] - bbox[0]
                x = (1920 - text_width) // 2
                draw.text((x, 100), slide_title, fill='black', font=font)
            
            # Draw content
            y_offset = 250
            for content in slide_content[:5]:  # Limit to 5 content items
                try:
                    content_font = ImageFont.truetype("arial.ttf", 24)
                except:
                    content_font = ImageFont.load_default()
                
                # Wrap long text
                if len(content) > 80:
                    content = content[:77] + "..."
                
                draw.text((100, y_offset), content, fill='black', font=content_font)
                y_offset += 60
            
            # Save image
            output_path = output_dir / f"slide_{i:02d}.png"
            img.save(output_path, 'PNG')
            print(f"Saved: {output_path.name}")
        
        print(f"\nConversion complete! Images saved to: {output_dir}")
        return True
        
    except Exception as e:
        print(f"Error converting presentation: {e}")
        return False

def main():
    if len(sys.argv) < 2:
        print("Usage: python convert_pptx_to_png.py <pptx_file> [output_directory]")
        print("Example: python convert_pptx_to_png.py presentation.pptx")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None
    
    success = convert_pptx_to_images(pptx_file, output_dir)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()